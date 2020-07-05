import os
from datetime import datetime
import xml.dom.minidom

currTime = datetime.now().strftime("%Y-%m-%d_%H")
#print(currTime)

# from datetime import datetime
# currTime = datetime.now().strftime("%d/%m/%Y %H:%M:%S")

#path = "C:\AAA_github\programming"

# def listdir_fullpath(d):
#     return [os.path.join(d, f) for f in os.listdir(d)]


# for f in listdir_fullpath(path):
#     print(f)




#    for name in dirs:
#       print(os.path.join(root, name))

# for top, dirs, files in os.walk(path):
#     for nm in files:       
#         print os.path.join(top, nm)

# print(os.path)

# for dirname, dirnames, filenames in os.walk(path):
#     # print path to all subdirectories first.
#     for subdirname in dirnames:
#         print(os.path.join(dirname, subdirname))

#     # print path to all filenames.
#     for filename in filenames:
#         print(os.path.join(dirname, filename))

#     # Advanced usage:
#     # editing the 'dirnames' list will stop os.walk() from recursing into there.
#     if '.git' in dirnames:
#         # don't go into any .git directories.
#         dirnames.remove('.git')


path = "C:\AAA files"
print("testing new ")
s = os.path.basename(__file__).split(".")[0]
print(s)
# for root, dirs, files in os.walk(path): # , topdown=False
#     for name in files:
#         print(os.path.join(root, name))


# x = [os.path.join(r,file) for r,d,f in os.walk(path) for file in f]

# for i in x:
#     print(i)




#======================================


import os
from datetime import datetime
import xml.dom.minidom


currTime = datetime.now().strftime("%Y-%m-%d_%H")
#print(currTime)

path = "C:\AAA files"
print("testing new ")
s = os.path.basename(__file__).split(".")[0]
print(s)


#===============================
# import PyPDF2
# filepath = "C:/Users/dalec/Desktop/3rd year/1st semester/CA357 User Interface and Completion/2_Human_Cognition_and_Perception/HumanPsychology.pdf"
# print(filepath)

# pdfFileObj = open(filepath,'rb')

# pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
# pages = pdfReader.numPages # prints the number of pages
# print(pages) 

# for i in range(0,3):
#     pageObj = pdfReader.getPage(i)
#         # Printing Page Number
#     print("Page No: ", i)
#     print("")

#     # text = pageObj.extractText().split("  ")
#     text = pageObj.extractText().split("  ")
#     print(text)

#     for i in range(len(text)):
#         # Printing the line
#         # Lines are seprated using "\n"
#         print(text[i])#,end="\n\n")
    

# pdfFileObj.close()
#===============================

filepath = "C:/Users/dalec/Desktop/3rd year/2nd semester/IT Architecture/002_Uni_ITA_Workshop_Tools_of_the_Trade_Recap_LCC_.ppt"

from pptx import Presentation
from pptx.util import Inches

prs = Presentation(filepath)

text_runs = []
for slide in prs.slides:
    print(slide)
#   for shape in slide.shapes:
#     if not shape.has_text_frame:
#       continue
#   for paragraph in shape.text_frame.paragraphs:
#     for run in paragraph.runs:
#       text_runs.append(run.text)